"""Slideshow (PPTX) export for chapter PNM rosters."""

from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO
from typing import Iterable, Optional, Protocol
import asyncio
from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import httpx

DEFAULT_ACCENT_RGB = (10, 10, 10)        # near-black
CREAM_RGB = (250, 247, 240)
TEXT_RGB = (10, 10, 10)
MUTED_RGB = (92, 92, 92)
BORDER_RGB = (232, 227, 214)


def hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
    s = hex_str.lstrip("#")
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def accent_or_default(theme: dict | None) -> tuple[int, int, int]:
    if not theme:
        return DEFAULT_ACCENT_RGB
    if not theme.get("enabled"):
        return DEFAULT_ACCENT_RGB
    hex_ = theme.get("accent_hex")
    if not hex_:
        return DEFAULT_ACCENT_RGB
    return hex_to_rgb(hex_)


def initials_for_name(name: str) -> str:
    parts = [p for p in name.strip().split() if p]
    if not parts:
        return "?"
    if len(parts) == 1:
        return parts[0][0].upper()
    return (parts[0][0] + parts[-1][0]).upper()


PHOTO_W, PHOTO_H = 600, 750  # 4:5 portrait, source size before placement on slide


def render_initials_avatar(
    name: str,
    accent_rgb: tuple[int, int, int],
    size: tuple[int, int] = (PHOTO_W, PHOTO_H),
) -> bytes:
    bg = (
        int(CREAM_RGB[0] * 0.93 + accent_rgb[0] * 0.07),
        int(CREAM_RGB[1] * 0.93 + accent_rgb[1] * 0.07),
        int(CREAM_RGB[2] * 0.93 + accent_rgb[2] * 0.07),
    )
    img = Image.new("RGB", size, color=bg)
    draw = ImageDraw.Draw(img)
    text = initials_for_name(name)

    font_size = int(size[1] * 0.4)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", font_size)
    except OSError:
        font = ImageFont.load_default()

    bbox = draw.textbbox((0, 0), text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    x = (size[0] - tw) / 2 - bbox[0]
    y = (size[1] - th) / 2 - bbox[1]
    draw.text((x, y), text, fill=accent_rgb, font=font)

    out = BytesIO()
    img.save(out, format="JPEG", quality=88)
    return out.getvalue()


def prepare_photo_bytes(data: bytes) -> bytes:
    img = Image.open(BytesIO(data))
    img = ImageOps.exif_transpose(img).convert("RGB")

    # center-crop to 4:5
    w, h = img.size
    target_ratio = 4 / 5
    current_ratio = w / h
    if current_ratio > target_ratio:
        new_w = int(h * target_ratio)
        left = (w - new_w) // 2
        img = img.crop((left, 0, left + new_w, h))
    elif current_ratio < target_ratio:
        new_h = int(w / target_ratio)
        top = (h - new_h) // 2
        img = img.crop((0, top, w, top + new_h))

    # cap long edge at 1200
    max_edge = 1200
    if max(img.size) > max_edge:
        ratio = max_edge / max(img.size)
        img = img.resize((int(img.size[0] * ratio), int(img.size[1] * ratio)))

    out = BytesIO()
    img.save(out, format="JPEG", quality=85, optimize=True)
    return out.getvalue()


@dataclass
class PnmSlideData:
    id: str
    name: str
    year: str
    major: str
    hometown: str
    status: str
    vote_summary: dict
    photo_bytes: Optional[bytes] = None
    tags: list[str] = field(default_factory=list)
    gpa: Optional[float] = None
    latest_note: Optional[dict] = None  # {"author": str, "text": str}


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _rgb(c: tuple[int, int, int]) -> RGBColor:
    return RGBColor(c[0], c[1], c[2])


def _add_background(slide, fill_rgb: tuple[int, int, int]):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(fill_rgb)
    bg.shadow.inherit = False
    return bg


def _add_accent_bar(slide, accent_rgb):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(73152))  # ~8px
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(accent_rgb)


def _add_text(slide, left, top, width, height, text, *, size_pt, rgb, bold=False, italic=False, font="Calibri"):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(rgb)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tx


def _build_cover(pres, chapter, theme, total, exported_at):
    slide = pres.slides.add_slide(pres.slide_layouts[6])  # blank
    _add_background(slide, CREAM_RGB)
    _add_accent_bar(slide, accent_or_default(theme))
    _add_text(slide, Inches(1), Inches(2.6), Inches(11.3), Inches(1.5),
              chapter["name"], size_pt=54, rgb=TEXT_RGB, bold=True, font="Georgia")
    _add_text(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.6),
              f"Rush — PNM Roster · {exported_at}", size_pt=18, rgb=MUTED_RGB)
    _add_text(slide, Inches(1), Inches(4.9), Inches(11.3), Inches(0.6),
              f"{total} PNMs", size_pt=14, rgb=MUTED_RGB)


def _build_section_divider(pres, theme, label, count):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    accent = accent_or_default(theme)
    soft = (
        int(CREAM_RGB[0] * 0.92 + accent[0] * 0.08),
        int(CREAM_RGB[1] * 0.92 + accent[1] * 0.08),
        int(CREAM_RGB[2] * 0.92 + accent[2] * 0.08),
    )
    _add_background(slide, soft)
    _add_text(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(1.2),
              f"{label} — {count} PNMs", size_pt=44, rgb=TEXT_RGB, bold=True, font="Georgia")


def _build_pnm_slide(pres, pnm: PnmSlideData, chapter, theme, page_number):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    _add_background(slide, CREAM_RGB)
    accent = accent_or_default(theme)
    _add_accent_bar(slide, accent)

    photo_bytes = pnm.photo_bytes or render_initials_avatar(pnm.name, accent)
    photo_stream = BytesIO(photo_bytes)
    slide.shapes.add_picture(photo_stream, Inches(0.7), Inches(0.8), width=Inches(3.5), height=Inches(4.375))

    info_left = Inches(4.7)
    _add_text(slide, info_left, Inches(0.8), Inches(8), Inches(0.9),
              pnm.name, size_pt=32, rgb=TEXT_RGB, bold=True, font="Georgia")
    _add_text(slide, info_left, Inches(1.7), Inches(8), Inches(0.5),
              f"{pnm.year} · {pnm.major}", size_pt=16, rgb=MUTED_RGB)

    gpa_str = f"{pnm.gpa:.2f}" if pnm.gpa is not None else ""
    if gpa_str:
        _add_text(slide, info_left, Inches(2.2), Inches(8), Inches(0.5),
                  f"GPA {gpa_str} · {pnm.hometown}", size_pt=14, rgb=MUTED_RGB)
    else:
        _add_text(slide, info_left, Inches(2.2), Inches(8), Inches(0.5),
                  pnm.hometown, size_pt=14, rgb=MUTED_RGB)

    if pnm.tags:
        chip_text = "Tags: " + ", ".join(pnm.tags[:5])
        if len(pnm.tags) > 5:
            chip_text += f", +{len(pnm.tags) - 5}"
        _add_text(slide, info_left, Inches(3.0), Inches(8), Inches(0.5),
                  chip_text, size_pt=12, rgb=TEXT_RGB)

    v = pnm.vote_summary
    _add_text(slide, info_left, Inches(3.7), Inches(8), Inches(0.5),
              f"Voting: {v.get('up',0)} bids · {v.get('down',0)} pass · {v.get('star',0)} star",
              size_pt=12, rgb=TEXT_RGB)

    if pnm.latest_note:
        note = pnm.latest_note["text"][:180]
        author = pnm.latest_note.get("author", "")
        _add_text(slide, info_left, Inches(4.5), Inches(8), Inches(0.5),
                  "Latest note:", size_pt=11, rgb=MUTED_RGB, bold=True)
        _add_text(slide, info_left, Inches(4.9), Inches(8), Inches(1.5),
                  f'"{note}"', size_pt=12, rgb=TEXT_RGB, italic=True)
        if author:
            _add_text(slide, info_left, Inches(6.0), Inches(8), Inches(0.4),
                      f"— {author}", size_pt=10, rgb=MUTED_RGB)

    _add_text(slide, Inches(10.8), Inches(7.1), Inches(2.3), Inches(0.3),
              f"{chapter['name']} · {page_number}", size_pt=9, rgb=MUTED_RGB)


def _build_closing(pres, exported_at):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    _add_background(slide, CREAM_RGB)
    _add_text(slide, Inches(1), Inches(3.4), Inches(11.3), Inches(0.6),
              "Exported from RushRank", size_pt=22, rgb=TEXT_RGB, bold=True, font="Georgia")
    _add_text(slide, Inches(1), Inches(4.1), Inches(11.3), Inches(0.4),
              exported_at, size_pt=12, rgb=MUTED_RGB)


_STATUS_LABEL = {
    "active": "Active",
    "round_1": "Round 1",
    "round_2": "Round 2",
    "round_3": "Round 3",
    "final": "Final",
    "bid": "Extended Bids",
    "cut": "Cut",
}


def build_deck(
    pnms: list[PnmSlideData],
    chapter: dict,
    theme: dict | None,
    exported_at: str,
) -> bytes:
    pres = Presentation()
    pres.slide_width = SLIDE_W
    pres.slide_height = SLIDE_H

    _build_cover(pres, chapter, theme, total=len(pnms), exported_at=exported_at)

    # group consecutive same-status into sections; insert divider when status changes
    page = 1
    last_status = None
    grouped: dict[str, int] = {}
    for p in pnms:
        grouped[p.status] = grouped.get(p.status, 0) + 1

    distinct_statuses = list(grouped.keys())
    insert_dividers = len(distinct_statuses) > 1

    for p in pnms:
        if insert_dividers and p.status != last_status:
            label = _STATUS_LABEL.get(p.status, p.status.title())
            _build_section_divider(pres, theme, label, grouped[p.status])
            last_status = p.status
        _build_pnm_slide(pres, p, chapter, theme, page_number=page)
        page += 1

    _build_closing(pres, exported_at)

    out = BytesIO()
    pres.save(out)
    return out.getvalue()


class _StorageClient(Protocol):
    async def fetch(self, path: str) -> bytes: ...


class SlideshowService:
    MAX_CONCURRENT_PHOTO_FETCH = 10

    def __init__(self, storage):
        self.storage = storage

    async def _fetch_one(self, url: str | None, sem: asyncio.Semaphore) -> bytes | None:
        if not url:
            return None
        async with sem:
            try:
                raw = await self.storage.fetch(url)
                return prepare_photo_bytes(raw)
            except Exception:
                return None

    async def build_pnm_deck(
        self,
        pnm_rows: list[dict],
        chapter: dict,
        theme: dict | None,
        exported_at: str,
    ) -> bytes:
        sem = asyncio.Semaphore(self.MAX_CONCURRENT_PHOTO_FETCH)
        photo_results = await asyncio.gather(
            *(self._fetch_one(row.get("photo_url"), sem) for row in pnm_rows)
        )

        pnms: list[PnmSlideData] = []
        for row, photo in zip(pnm_rows, photo_results):
            pnms.append(PnmSlideData(
                id=str(row["id"]),
                name=row["name"],
                year=row.get("year", ""),
                major=row.get("major", ""),
                gpa=float(row.get("gpa") or 0.0) if row.get("gpa") else None,
                hometown=row.get("hometown", ""),
                status=row.get("status", "active"),
                vote_summary=row.get("vote_summary") or {"up": 0, "down": 0, "star": 0},
                photo_bytes=photo,
                tags=list(row.get("tags") or []),
                latest_note=row.get("latest_note"),
            ))

        return build_deck(pnms, chapter=chapter, theme=theme, exported_at=exported_at)


class HttpPhotoFetcher:
    """Fetches photo bytes from arbitrary http(s) URLs.

    `fetch(url)` returns bytes or raises. Compatible with SlideshowService's
    storage Protocol (any object with async `fetch(path) -> bytes`).
    """
    def __init__(self, timeout_s: float = 8.0):
        self.timeout_s = timeout_s

    async def fetch(self, url: str) -> bytes:
        if not url:
            raise FileNotFoundError("empty url")
        async with httpx.AsyncClient(timeout=self.timeout_s) as client:
            r = await client.get(url)
            r.raise_for_status()
            return r.content
