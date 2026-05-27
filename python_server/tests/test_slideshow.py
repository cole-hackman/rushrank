import pytest
from io import BytesIO
from PIL import Image
from pptx import Presentation
from python_server.slideshow import (
    hex_to_rgb,
    accent_or_default,
    initials_for_name,
    render_initials_avatar,
    prepare_photo_bytes,
    build_deck,
    PnmSlideData,
)


def test_hex_to_rgb_parses_uppercase():
    assert hex_to_rgb("#0033A0") == (0, 51, 160)


def test_hex_to_rgb_parses_lowercase():
    assert hex_to_rgb("#ffc0cb") == (255, 192, 203)


def test_accent_or_default_returns_default_when_disabled():
    theme = {"enabled": False, "accent_hex": "#0033A0", "source": "auto"}
    assert accent_or_default(theme) == (10, 10, 10)


def test_accent_or_default_returns_default_when_no_hex():
    theme = {"enabled": True, "accent_hex": None, "source": "auto"}
    assert accent_or_default(theme) == (10, 10, 10)


def test_accent_or_default_returns_accent_when_enabled():
    theme = {"enabled": True, "accent_hex": "#0033A0", "source": "manual"}
    assert accent_or_default(theme) == (0, 51, 160)


def test_initials_single_word():
    assert initials_for_name("Madonna") == "M"


def test_initials_two_words():
    assert initials_for_name("Marcus Chen") == "MC"


def test_initials_three_plus_words():
    assert initials_for_name("Mary Jane Watson") == "MW"


def test_initials_handles_empty():
    assert initials_for_name("") == "?"


def test_render_initials_avatar_returns_jpeg_bytes():
    data = render_initials_avatar("Marcus Chen", accent_rgb=(0, 51, 160))
    img = Image.open(BytesIO(data))
    assert img.format == "JPEG"
    assert img.size == (600, 750)


def test_prepare_photo_bytes_crops_to_4x5():
    # synthetic landscape jpeg
    src = Image.new("RGB", (1200, 800), color=(120, 120, 120))
    buf = BytesIO()
    src.save(buf, format="JPEG")
    buf.seek(0)
    out = prepare_photo_bytes(buf.getvalue())
    img = Image.open(BytesIO(out))
    assert img.format == "JPEG"
    w, h = img.size
    assert abs((w / h) - (4 / 5)) < 0.01
    assert max(w, h) <= 1200


def test_prepare_photo_bytes_handles_portrait_taller_than_4x5():
    src = Image.new("RGB", (600, 1200), color=(120, 120, 120))
    buf = BytesIO()
    src.save(buf, format="JPEG")
    buf.seek(0)
    out = prepare_photo_bytes(buf.getvalue())
    img = Image.open(BytesIO(out))
    w, h = img.size
    assert abs((w / h) - (4 / 5)) < 0.01


SAMPLE_PNM = lambda i, **kw: PnmSlideData(
    id=f"pnm-{i}",
    name=kw.get("name", f"PNM {i}"),
    year=kw.get("year", "Freshman"),
    major=kw.get("major", "Economics"),
    gpa=kw.get("gpa", 3.7),
    hometown=kw.get("hometown", "Boston, MA"),
    photo_bytes=kw.get("photo_bytes", None),
    tags=kw.get("tags", []),
    status=kw.get("status", "active"),
    vote_summary=kw.get("vote_summary", {"up": 5, "down": 1, "star": 0}),
    latest_note=kw.get("latest_note", None),
)

CHAPTER_META = {"name": "Sigma Chi at BC", "fraternity": "Sigma Chi"}
THEME_ON = {"enabled": True, "accent_hex": "#0033A0", "source": "manual"}
THEME_OFF = {"enabled": False, "accent_hex": None, "source": "auto"}


def test_build_deck_has_cover_and_closing():
    pnms = [SAMPLE_PNM(i) for i in range(3)]
    data = build_deck(pnms, chapter=CHAPTER_META, theme=THEME_OFF, exported_at="2026-05-27")
    pres = Presentation(BytesIO(data))
    # cover + 3 pnms + closing = 5 (no section dividers since all same status)
    assert len(pres.slides) == 5


def test_build_deck_inserts_section_divider_between_statuses():
    pnms = [
        SAMPLE_PNM(1, status="round_1"),
        SAMPLE_PNM(2, status="round_1"),
        SAMPLE_PNM(3, status="round_2"),
    ]
    data = build_deck(pnms, chapter=CHAPTER_META, theme=THEME_OFF, exported_at="2026-05-27")
    pres = Presentation(BytesIO(data))
    # cover + round_1 divider + 2 pnms + round_2 divider + 1 pnm + closing = 7
    assert len(pres.slides) == 7


def test_build_deck_uses_accent_when_theme_enabled():
    pnms = [SAMPLE_PNM(1)]
    data = build_deck(pnms, chapter=CHAPTER_META, theme=THEME_ON, exported_at="2026-05-27")
    pres = Presentation(BytesIO(data))
    # find the accent bar shape on the PNM slide (slide index 1, after cover)
    accent_shapes = [s for s in pres.slides[1].shapes if s.shape_type == 1]  # rectangle
    # at least one rectangle in the accent color (0, 51, 160)
    found = False
    for s in accent_shapes:
        if s.fill.fore_color.rgb == (0, 51, 160):
            found = True
            break
    assert found, "no accent-colored shape on PNM slide"


def test_build_deck_missing_photo_uses_avatar():
    pnms = [SAMPLE_PNM(1, photo_bytes=None, name="Marcus Chen")]
    data = build_deck(pnms, chapter=CHAPTER_META, theme=THEME_OFF, exported_at="2026-05-27")
    pres = Presentation(BytesIO(data))
    # slide 1 (cover=0, first PNM=1)
    pics = [s for s in pres.slides[1].shapes if s.shape_type == 13]
    assert len(pics) == 1  # avatar still injected as a picture


def test_build_deck_pnm_slide_contains_name_and_metadata():
    pnms = [SAMPLE_PNM(1, name="Marcus Chen", year="Freshman", major="CS", gpa=3.91)]
    data = build_deck(pnms, chapter=CHAPTER_META, theme=THEME_OFF, exported_at="2026-05-27")
    pres = Presentation(BytesIO(data))
    texts = []
    for shape in pres.slides[1].shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    texts.append(r.text)
    full = " ".join(texts)
    assert "Marcus Chen" in full
    assert "Freshman" in full
    assert "CS" in full
    assert "3.91" in full
